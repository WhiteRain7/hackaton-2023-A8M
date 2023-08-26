import collections  # noqa: F401
import collections.abc  # noqa: F401
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN

from ..schemas.base_schema import CreatePresentation, MarketType
from ..config import settings
from fastapi.responses import FileResponse
import anyio
from starlette.background import BackgroundTask
from time import time
from pptx.shapes.placeholder import TablePlaceholder

from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
from pptx.dml.color import RGBColor

from io import BytesIO
import matplotlib.dates as mdates
import matplotlib.pyplot as plt
from matplotlib.figure import Figure

from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_TICK_MARK
from pptx.enum.chart import XL_LABEL_POSITION

from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
import random
import httpx

import io
import base64
from PIL import Image, PngImagePlugin
from ..hugchat.api import call_api

class BaseSevice:
    def generate_template(self) -> Presentation:
        pres_temp = random.randint(1,10)
        prs = Presentation(str(settings.templates_path / f'{pres_temp}.pptx'))
        return prs

    def get_top_min(self, slide) -> int:
        _title = slide.shapes.title
        return _title.top + _title.height


    def generate_filename(self) -> Path:
        return settings.media_path / f"{time()}.pptx"
    
    def create_table(self, slide, _rows, _cols):
        left_inch = Inches(2.65)
        top_inch = self.get_top_min(slide) + Inches(0.5)
        width_inch = Inches(8.0)
        height_inch = Inches(3)
        table = slide.shapes.add_table(
            rows=_rows, cols=_cols,
            left=left_inch, top=top_inch, width=width_inch, height=height_inch
        ).table

        table.first_col = False
        table.first_row = True
                        
        return table
    
    def create_slide(self,  prs, slide_layout, title):
        title_slide_layout = prs.slide_layouts[slide_layout]
        slide = prs.slides.add_slide(title_slide_layout)
        _title = slide.shapes.title
        _title.text = title
        return slide

    def generate_title_slide(self, prs: Presentation, data: CreatePresentation) -> None:
        """Генерация первого слайда - Титульник"""
        slide = self.create_slide(prs, 0, data.project_name)
        _subtitle = slide.placeholders[1]
        _subtitle.text = data.short_description

    def generate_table_slide(self, prs: Presentation, data: CreatePresentation, title: str, content_func):
        slide = self.create_slide(prs, 5, title)

        max_content_length = max(len(content_func(x)) for x in data.problem)
        _rows = max_content_length + 1
        _cols = len(data.problem)

        table = self.create_table(slide, _rows, _cols)

        for i in range(0, _cols):
            cell = table.cell(0, i)
            cell.text = data.problem[i].target_audience
            for j, content in enumerate(content_func(data.problem[i])):
                cell = table.cell(j + 1, i)
                cell.text = content

    def generate_problems_slide(self, prs: Presentation, data: CreatePresentation):
        """Генерация второго слайда - Проблемы"""
        self.generate_table_slide(prs, data, "Проблемы", lambda x: x.issue)

    def generate_solution_slide(self, prs: Presentation, data: CreatePresentation):
        """Генерация четвёртого слайда - Решение"""
        self.generate_table_slide(prs, data, "Решение", lambda x: x.solution)

    async def generate_description_slide(self, prs: Presentation, data: CreatePresentation) -> None:
        """Генерация третьего слайда - Описание"""
        slide = self.create_slide(prs, 3, "Описание")
        description = slide.placeholders[1]

        prompt = self.generate_prompt(data.description)
        image_stream = await self.generate_image(f'description image high quality {prompt}')
        left = Inches(7.3)
        top = self.get_top_min(slide) + Inches(0.5)
        width = Inches(4.5)
        height = Inches(4.5)
        
        slide.shapes.add_picture(image_stream, left, top, width, height)
        description.text = data.description

    def generate_members_slide(self, prs: Presentation, data: CreatePresentation):
        """Генерация 10 слайда - Команда"""
        slide = self.create_slide(prs, 5, "Команда")
        _rows = 2
        _cols = len(data.members)
        table = self.create_table(slide, _rows, _cols)
        for i in range(0, _cols):
            table.cell(0, i).text = data.members[i].full_name
            table.cell(1, i).text = data.members[i].proffesion

    def generate_investors_slide(self, prs: Presentation, data: CreatePresentation):
        """Генерация 11 слайда - Инвесторы"""
        slide = self.create_slide(prs, 1, "Инвесторы")
        description = slide.placeholders[1]
        description.text = '\n'.join(item.name for item in data.investors) + '\n'

    def generate_opponents_slide(self, prs: Presentation, data: CreatePresentation):
        """Генерация 6 слайда - Конкуренты"""
        slide = self.create_slide(prs, 1, "Конкуренты")
        description = slide.placeholders[1]
        description.text = '\n'.join(data.opponents) + '\n'

    def generate_market_slide(self, prs: Presentation, data: CreatePresentation):
        """Генерация 5 слайда - Рынок"""
        slide = self.create_slide(prs, 5, "Рынок")
        left = Inches(1.8)
        top = self.get_top_min(slide) + Inches(0.7)
        width = Inches(3)
        height = Inches(3)
        shapes = slide.shapes
        for i, market_unit in enumerate(data.market):
            left_offset = i * Inches(3.25)
            shape = shapes.add_shape(
                MSO_SHAPE.OVAL, left + left_offset, top, width, height
            )
            
            head = shape.text_frame.add_paragraph()
            head.text = market_unit.type.value
            head.font.bold = True
            head.font.size = Pt(32)
            head.alignment = PP_ALIGN.CENTER
            
            volume = shape.text_frame.add_paragraph()
            volume.text = market_unit.volume
            volume.font.bold = True
            volume.font.size = Pt(24)
            volume.alignment = PP_ALIGN.CENTER

            name = shape.text_frame.add_paragraph()
            name.text = market_unit.name
            name.font.bold = False
            name.alignment = PP_ALIGN.CENTER
            name.margin_top = Pt(8)
    
    def generate_investing_rounds_slide(self, prs: Presentation, data: CreatePresentation):
        """Генерация 12 слайда - Инвестиционный раунд"""
        investing_rounds = data.investing_rounds
        if len(investing_rounds) > 1:
            raise Exception("Not implemented")

        stage = investing_rounds[0].stage
        amount = investing_rounds[0].amount
        slide = self.create_slide(prs, 5, "")
        
        title_text = f"За {stage.value} будет привлечено {amount} рублей"
        if investing_rounds[0].fraction:
            title_text += f" за долю в компании {investing_rounds[0].fraction}"
        _title = slide.shapes.title
        _title.text = title_text
        
        spending = investing_rounds[0].spending
        labels = [item.name for item in spending]
        sizes = [int(item.percent.strip('%')) for item in spending]
        
        fig1, ax1 = plt.subplots(figsize=(8, 4))
        ax1.pie(sizes, labels=labels, autopct='%1.1f%%', startangle=90)
        ax1.axis('equal')

        image_stream = BytesIO()
        plt.savefig(image_stream, format='png')
        plt.close(fig1)
        
        left = Inches(2.5)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(4)
        
        pic = slide.shapes.add_picture(image_stream, left, top, width, height)
        
        image_stream.seek(0)
    
    def generate_roadmap_slide(self, prs: Presentation, data: CreatePresentation):
        """Генерация 13 слайда - Дорожная карта"""
        slide = self.create_slide(prs, 5, "Дорожная карта")
        fig, ax = plt.subplots(figsize=(9, 4))
        steps = data.roadmap
        for i, step in enumerate(steps):
            ax.barh(i, (step.end_date - step.start_date).days, left=step.start_date, height=0.5, align='center')
        ax.set_yticks(range(len(steps)))
        ax.set_yticklabels([step.name for step in steps])
        ax.xaxis_date()
        # FIXME: дату на основе месяца/года
        ax.xaxis.set_major_locator(mdates.DayLocator(interval=30))
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%Y-%m-%d'))

        image_stream = BytesIO()
        plt.savefig(image_stream, format='png')
        plt.close()

        left = Inches(1.9)
        top = self.get_top_min(slide) + Inches(0.5)
        width = Inches(9.5)
        height = Inches(4)
        pic = slide.shapes.add_picture(image_stream, left, top, width, height)
    
    def generate_contacts_slide(self, prs: Presentation, data: CreatePresentation):
        """Генерация 14 слайда - Контакты"""
        slide = self.create_slide(prs, 1, "Контакты")
        description = slide.placeholders[1]
        description.text = '\n'.join(item.value for item in data.contacts) + '\n'
    
    def generate_business_units_slide(self, prs: Presentation, data: CreatePresentation):
        """Генерация 7 слайда - Бизнес-модель"""
        slide = self.create_slide(prs, 5, "Бизнес-модель")

        _rows = 3
        _cols = len(data.business_units)

        table = self.create_table(slide, _rows, _cols)

        for i in range(0, _cols):
            table.cell(0, i).text = data.business_units[i].name
            table.cell(1, i).text = data.business_units[i].income
            table.cell(2, i).text = data.business_units[i].revenue_share

    def generate_tracktion_slide(self, prs: Presentation, data: CreatePresentation):
        """Генерация 8 слайда - Трекшен"""
        slide = self.create_slide(prs, 5, "Трекшен")
        
        left = Inches(3)
        top = self.get_top_min(slide) + Inches(0.5)
        width = Inches(10)
        height = Inches(4)
        
        tracktion = data.tracktion
        years = [str(item.year) for item in tracktion]
        revenue = [int(item.revenue) for item in tracktion]
        capitalization = [int(item.capitalization) for item in tracktion]
        
        fig, ax1 = plt.subplots(figsize=(10, 4))
        ax2 = ax1.twinx()
        
        ax2.plot(years, revenue, label='Выручка', color='tab:orange', linewidth=3)
        ax2.set_xlabel('Год')
        ax2.set_ylabel('Выручка', color='tab:orange')
        ax2.tick_params(axis='y', labelcolor='tab:orange')

        ax1.bar(years, capitalization, label='Капитализация', color='tab:blue')
        ax1.set_ylabel('Капитализация', color='tab:blue')
        ax1.tick_params(axis='y', labelcolor='tab:blue')
        
        lines, labels = ax1.get_legend_handles_labels()
        lines2, labels2 = ax2.get_legend_handles_labels()
        ax2.legend(lines + lines2, labels + labels2, loc='upper left')
        
        plt.tight_layout()

        image_stream = BytesIO()
        plt.savefig(image_stream, format='png')
        plt.close()
        image_stream.seek(0)
        
        left = Inches(1.8)
        top = self.get_top_min(slide) + Inches(0.5)
        pic = slide.shapes.add_picture(image_stream, left, top, width, height)
    
    def generate_finance_slide(self, prs: Presentation, data: CreatePresentation):
        """Генерация 9 слайда - Финансы"""
        slide = self.create_slide(prs, 1, "Финансы")
        description = slide.placeholders[1]
        description.text = f"\nВыручка: {data.revenue} руб\nЧисло клиентов: {data.clients_count}\nChurn Rate: {data.churn_rate}%\nAPRU = {int(data.revenue) / int(data.clients_count)} руб\nLT = {1 / (int(data.churn_rate) / 100)} лет\nLTV = {(int(data.revenue) // int(data.clients_count)) * (1 // (int(data.churn_rate) / 100))}\n"

    async def create_presentation(self, data: CreatePresentation):
        prs = self.generate_template()
        # TODO: цвета + шрифты
        file = self.generate_filename()

        self.generate_title_slide(prs, data)
        self.generate_problems_slide(prs, data)
        await self.generate_description_slide(prs, data)
        self.generate_solution_slide(prs, data)
        self.generate_market_slide(prs, data)
        if (len(data.opponents)):
            self.generate_opponents_slide(prs, data)
        self.generate_business_units_slide(prs, data)
        self.generate_tracktion_slide(prs, data)
        self.generate_finance_slide(prs, data)
        self.generate_members_slide(prs, data)
        if (len(data.investors)):
            self.generate_investors_slide(prs, data)
        self.generate_investing_rounds_slide(prs, data)
        self.generate_roadmap_slide(prs, data)
        self.generate_contacts_slide(prs, data)

        prs.save(str(file))
        return FileResponse(
            file, headers={"Content-Disposition": f'attachment; filename="{file.name}"'}, background=BackgroundTask(self.delete_file, file)
        )  # background=BackgroundTask(self.delete_file, file)

    def generate_prompt(self, value: str) -> str:
        return call_api(f"create an 5 base key words list based on next text separated with comma on english:{value}").split(":")[-1].replace("\n", ", ").replace(" ", ", ")

    async def generate_image(self, prompt: str) -> io.BytesIO:
        async with httpx.AsyncClient(timeout=1000) as client:
            try:
                r = await client.post(f"{settings.stable_diffusion_url}sdapi/v1/txt2img", json={"prompt": prompt, "steps":30, "width": 512,"height": 512,})
            except:
                with open(settings.templates_path / "zaglushka.png", "rb") as file:
                    buffer = io.BytesIO(file.read())
                return buffer
            else:
                r = r.json()
                for i in r['images']:
                    image = io.BytesIO(base64.b64decode(i.split(",",1)[0]))
                    return image

    async def delete_file(self, file):
        await anyio.run_sync_in_worker_thread(file.unlink)


# https://github.com/Soulter/hugging-chat-api
